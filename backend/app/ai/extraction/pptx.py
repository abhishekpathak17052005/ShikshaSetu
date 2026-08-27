"""PPTX presentation extraction."""
from typing import List, Tuple

from pptx import Presentation


class PPTXExtractor:
    """Extract text from PPTX presentations."""

    @staticmethod
    def extract(file_path: str) -> Tuple[str, List[dict]]:
        """
        Extract text from a PPTX file.

        Args:
            file_path: Path to the PPTX file.

        Returns:
            Tuple of:
            - Full extracted text
            - List of slide metadata with text

        Raises:
            Exception: If PPTX cannot be read or has no extractable text.
        """
        try:
            prs = Presentation(file_path)
            
            if len(prs.slides) == 0:
                raise Exception("PPTX contains no slides")
            
            full_text = []
            slides_metadata = []
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                    
                    # Extract table text if present
                    if shape.has_table:
                        table = shape.table
                        table_rows = []
                        for row in table.rows:
                            row_cells = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                if cell_text:
                                    row_cells.append(cell_text)
                            if row_cells:
                                table_rows.append(" | ".join(row_cells))
                        
                        if table_rows:
                            slide_text.append("\n".join(table_rows))
                
                # Extract speaker notes if available
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_text.append(f"[Notes: {notes_text}]")
                
                if slide_text:
                    combined_slide_text = "\n".join(slide_text)
                    full_text.append(combined_slide_text)
                    slides_metadata.append({
                        "slide": slide_num,
                        "text": combined_slide_text
                    })
            
            if not full_text:
                raise Exception("PPTX contains no extractable text")
            
            return "\n".join(full_text), slides_metadata
        
        except Exception as e:
            # Re-raise with context
            raise Exception(f"PPTX extraction failed: {str(e)}")
