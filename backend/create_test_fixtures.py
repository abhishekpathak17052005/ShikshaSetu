"""Create realistic test fixtures for Phase 6 verification."""
import os
from pathlib import Path

# Create fixtures directory
fixtures_dir = Path("tests/fixtures")
fixtures_dir.mkdir(exist_ok=True)

# Create a simple PDF with PyPDF2
print("Creating sample SQL PDF...")
try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    
    pdf_path = fixtures_dir / "sample_sql.pdf"
    c = canvas.Canvas(str(pdf_path), pagesize=letter)
    
    # Page 1
    c.setFont("Helvetica-Bold", 16)
    c.drawString(50, 750, "SQL Tutorial: Database Queries")
    
    c.setFont("Helvetica", 12)
    c.drawString(50, 720, "Page 1: SQL SELECT Statement")
    c.drawString(50, 700, "The SELECT statement retrieves data from one or more database tables.")
    c.drawString(50, 680, "Syntax: SELECT column1, column2 FROM table_name WHERE condition;")
    
    c.drawString(50, 650, "Page 1: WHERE Clause")
    c.drawString(50, 630, "A WHERE clause filters rows according to a specified condition.")
    c.drawString(50, 610, "Only records that meet the condition are included in the result.")
    
    # Page 2
    c.showPage()
    c.setFont("Helvetica", 12)
    c.drawString(50, 750, "Page 2: JOIN Operations")
    c.drawString(50, 730, "An INNER JOIN returns rows when matching values exist in both tables.")
    c.drawString(50, 710, "Syntax: SELECT * FROM table1 INNER JOIN table2 ON table1.id = table2.id;")
    c.drawString(50, 680, "A LEFT JOIN returns all rows from the left table and matching rows from right table.")
    
    c.save()
    print(f"✓ Created {pdf_path} ({pdf_path.stat().st_size} bytes)")
except Exception as e:
    print(f"✗ Could not create PDF with reportlab: {e}")
    print("  Creating minimal PDF manually...")
    pdf_path = fixtures_dir / "sample_sql.pdf"
    with open(pdf_path, "wb") as f:
        f.write(b"""%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj
2 0 obj
<< /Type /Pages /Kids [3 0 R] /Count 1 >>
endobj
3 0 obj
<< /Type /Page /Parent 2 0 R /Resources 4 0 R /MediaBox [0 0 612 792] /Contents 5 0 R >>
endobj
4 0 obj
<< /Font << /F1 << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> >> >>
endobj
5 0 obj
<< /Length 200 >>
stream
BT
/F1 12 Tf
50 700 Td
(SQL Tutorial) Tj
0 -20 Td
(SELECT retrieves data from database tables.) Tj
0 -20 Td
(WHERE clause filters rows by condition.) Tj
0 -20 Td
(INNER JOIN returns matching rows from both tables.) Tj
ET
endstream
endobj
xref
0 6
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000203 00000 n
0000000286 00000 n
trailer
<< /Size 6 /Root 1 0 R >>
startxref
536
%%EOF""")
    print(f"✓ Created {pdf_path} ({pdf_path.stat().st_size} bytes)")

# Create a DOCX file
print("\nCreating sample Python DOCX...")
try:
    from docx import Document
    from docx.shared import Pt
    
    docx_path = fixtures_dir / "sample_python.docx"
    doc = Document()
    
    doc.add_heading("Python Programming Guide", 0)
    doc.add_paragraph("Python is a high-level, interpreted programming language.")
    doc.add_paragraph("It emphasizes code readability and simplicity.")
    
    doc.add_heading("Variables and Data Types", level=1)
    doc.add_paragraph("Python variables do not require explicit type declaration.")
    doc.add_paragraph("Common data types: int, float, str, list, dict, tuple, set.")
    
    doc.add_heading("Functions", level=1)
    doc.add_paragraph("Functions are defined with the def keyword.")
    doc.add_paragraph("Example: def greet(name): return f'Hello, {name}'")
    
    doc.add_heading("Control Flow", level=1)
    doc.add_paragraph("If-else statements control program flow based on conditions.")
    doc.add_paragraph("Loops (for, while) iterate over sequences or conditions.")
    
    doc.save(docx_path)
    print(f"✓ Created {docx_path} ({docx_path.stat().st_size} bytes)")
except Exception as e:
    print(f"✗ Error creating DOCX: {e}")

# Create a PPTX file
print("\nCreating sample Introduction PPTX...")
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    pptx_path = fixtures_dir / "sample_intro.pptx"
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "Introduction to Cloud Computing"
    title_frame.paragraphs[0].font.size = Pt(44)
    
    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "What is Cloud Computing?"
    title_frame.paragraphs[0].font.size = Pt(32)
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = "Cloud computing provides on-demand access to computing resources over the internet."
    content_frame.word_wrap = True
    
    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Cloud Service Models"
    title_frame.paragraphs[0].font.size = Pt(32)
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
    content_frame = content_box.text_frame
    p1 = content_frame.paragraphs[0]
    p1.text = "IaaS: Infrastructure as a Service"
    p2 = content_frame.add_paragraph()
    p2.text = "PaaS: Platform as a Service"
    p3 = content_frame.add_paragraph()
    p3.text = "SaaS: Software as a Service"
    
    prs.save(pptx_path)
    print(f"✓ Created {pptx_path} ({pptx_path.stat().st_size} bytes)")
except Exception as e:
    print(f"✗ Error creating PPTX: {e}")

print("\n✓ Test fixtures created successfully!")
